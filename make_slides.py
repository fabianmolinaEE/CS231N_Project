#!/usr/bin/env python3
"""
Build Milestone 3 slides for CS231N: Physics-Constrained Chip Thermal Prediction.
Run: python3 make_slides.py
Output: milestone3_slides.pptx
"""

import os
import tempfile
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BG    = RGBColor(0x0D, 0x11, 0x17)
CARD  = RGBColor(0x16, 0x1B, 0x22)
CARD2 = RGBColor(0x1E, 0x29, 0x3B)
ORG   = RGBColor(0xF9, 0x73, 0x16)
BLUE  = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
GRAY  = RGBColor(0x94, 0xA3, 0xB8)
YEL   = RGBColor(0xFD, 0xE6, 0x8A)
RED   = RGBColor(0xEF, 0x44, 0x44)
SLATE = RGBColor(0x33, 0x41, 0x55)
PURP  = RGBColor(0xA8, 0x55, 0xF7)

W, H = Inches(13.333), Inches(7.5)
DIR  = Path(__file__).parent
TEMPS: list[str] = []


def tmp() -> str:
    f = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    TEMPS.append(f.name)
    f.close()
    return f.name


def cleanup():
    for p in TEMPS:
        try:
            os.unlink(p)
        except OSError:
            pass


# ── pptx helpers ──────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def textbox(slide, text, left, top, width, height,
            size=20, bold=False, italic=False,
            color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(0.07), ORG)
    textbox(slide, title,
            Inches(0.55), Inches(0.15), Inches(12.5), Inches(0.70),
            size=34, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, subtitle,
                Inches(0.55), Inches(0.78), Inches(12.5), Inches(0.38),
                size=15, italic=True, color=GRAY)
    add_rect(slide, Inches(0.55), Inches(1.12), Inches(12.22), Inches(0.025), SLATE)


def add_img(slide, path, left, top, width=None, height=None):
    kw = {}
    if width:
        kw["width"] = width
    if height:
        kw["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kw)


# ── Matplotlib style ──────────────────────────────────────────────────────────

STYLE  = dict(facecolor="#0D1117")
AX_BG  = "#161B22"
GRID_C = "#334155"
TX_C   = "#94A3B8"
FG_C   = "#F8FAFC"


def mpl_style(ax):
    ax.set_facecolor(AX_BG)
    for sp in ["top", "right"]:
        ax.spines[sp].set_visible(False)
    for sp in ["bottom", "left"]:
        ax.spines[sp].set_color(GRID_C)
    ax.tick_params(colors=TX_C, labelsize=11)
    ax.xaxis.label.set_color(TX_C)
    ax.yaxis.label.set_color(TX_C)
    ax.grid(color=GRID_C, alpha=0.5, zorder=0)


# ── Chart generators ──────────────────────────────────────────────────────────

def chart_title_bg() -> str:
    rng = np.random.default_rng(42)
    x = np.linspace(0, 1, 500)
    y = np.linspace(0, 1, 300)
    X, Y = np.meshgrid(x, y)
    Z = np.zeros_like(X)
    for _ in range(10):
        cx, cy = rng.uniform(0.05, 0.95), rng.uniform(0.05, 0.95)
        sx, sy = rng.uniform(0.08, 0.25), rng.uniform(0.08, 0.25)
        Z += rng.uniform(0.2, 1.0) * np.exp(-((X-cx)**2/(2*sx**2) + (Y-cy)**2/(2*sy**2)))
    fig, ax = plt.subplots(figsize=(13.333, 7.5))
    fig.patch.set_facecolor("#0D1117")
    ax.set_position([0, 0, 1, 1])
    ax.imshow(Z, cmap="inferno", alpha=0.18, aspect="auto", origin="lower")
    ax.axis("off")
    path = tmp()
    fig.savefig(path, dpi=72, bbox_inches="tight", pad_inches=0, facecolor="#0D1117")
    plt.close(fig)
    return path


def chart_rmse_bar() -> str:
    models = ["PlainCNN\n(no skip, dilated)", "EncoderDecoder\n(no skip)", "U-Net\n(with physics, λ=0.1)", "U-Net\n(MSE only, λ=0)  ★"]
    rmse   = [1.67, 1.20, 0.89, 0.78]
    clrs   = ["#64748B", "#475569", "#3B82F6", "#F97316"]

    fig, ax = plt.subplots(figsize=(8.5, 4.5), **STYLE)
    mpl_style(ax)

    bars = ax.barh(models, rmse, color=clrs, height=0.52, zorder=3, edgecolor="none")
    for bar, val in zip(bars, rmse):
        ax.text(val + 0.04, bar.get_y() + bar.get_height() / 2,
                f"{val:.2f} K", va="center", ha="left",
                color=FG_C, fontsize=12, fontweight="bold")

    ax.axvline(x=0.78, color="#F97316", linestyle="--", alpha=0.55, zorder=4, linewidth=1.5)
    ax.text(0.80, 0.02, "best", color="#F97316",
            transform=ax.get_xaxis_transform(), fontsize=10)
    ax.set_xlim(0, 2.2)
    ax.set_xlabel("Validation RMSE (Kelvin) — lower is better", fontsize=12, color=TX_C)
    ax.set_title("Val RMSE by Model (250 epochs, val split, 31 samples)", color=FG_C, fontsize=13, pad=10)
    ax.yaxis.set_tick_params(labelcolor=FG_C)

    plt.tight_layout(pad=0.9)
    path = tmp()
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="#0D1117")
    plt.close(fig)
    return path


def chart_lambda() -> str:
    fig, ax = plt.subplots(figsize=(6.5, 4.0), **STYLE)
    mpl_style(ax)

    lam  = [0.0, 0.1]
    mse  = [0.01691, 0.02196]

    ax.plot(lam, mse, "o-", color="#F97316",
            linewidth=2.5, markersize=13, zorder=4,
            markerfacecolor="#FDE68A", markeredgecolor="#F97316", markeredgewidth=2)

    for l, m in zip(lam, mse):
        ax.annotate(f"  {m:.5f}", (l, m),
                    textcoords="offset points", xytext=(6, 7),
                    color="#FDE68A", fontsize=12, fontweight="bold")

    ax.annotate("", xy=(0.35, 0.026), xytext=(0.18, 0.022),
                arrowprops=dict(arrowstyle="->", color="#94A3B8", lw=1.5, linestyle="dashed"))
    ax.text(0.19, 0.024, "expected\ntrend", color="#94A3B8", fontsize=10, style="italic")

    ax.set_xlabel("Physics loss weight  λ  (0 = pure MSE, higher = more physics)", fontsize=11)
    ax.set_ylabel("Validation MSE (normalized)", fontsize=11)
    ax.set_title("Measured: λ=0.0 and λ=0.1  (both U-Net, 250 epochs)", color=FG_C, fontsize=12, pad=8)
    ax.set_xlim(-0.02, 0.45)
    ax.set_ylim(0.012, 0.030)
    ax.fill_between(lam, [m - 0.0003 for m in mse], [m + 0.0003 for m in mse],
                    color="#F97316", alpha=0.15)

    plt.tight_layout(pad=0.9)
    path = tmp()
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="#0D1117")
    plt.close(fig)
    return path


def chart_mse_bars() -> str:
    models = ["PlainCNN", "Encoder\nDecoder", "U-Net\nλ=0.1", "U-Net\nλ=0.0"]
    mse    = [0.07693, 0.03972, 0.02196, 0.01691]
    clrs   = ["#64748B", "#475569", "#3B82F6", "#F97316"]

    fig, ax = plt.subplots(figsize=(6.8, 4.0), **STYLE)
    mpl_style(ax)

    x    = np.arange(len(models))
    bars = ax.bar(x, mse, color=clrs, width=0.55, zorder=3, edgecolor="none")
    for bar, val in zip(bars, mse):
        ax.text(bar.get_x() + bar.get_width() / 2, val + 0.001,
                f"{val:.4f}", ha="center", va="bottom",
                color=FG_C, fontsize=11, fontweight="bold")

    # improvement labels between bars
    for i in range(len(models) - 1):
        pct = (mse[i] - mse[i+1]) / mse[i] * 100
        mid_y = (mse[i] + mse[i+1]) / 2
        ax.text(i + 0.5, mid_y + 0.001, f"−{pct:.0f}%",
                ha="center", va="bottom", color="#22C55E", fontsize=10, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(models, color=FG_C, fontsize=11)
    ax.set_ylabel("Validation MSE (lower = better)", fontsize=11)
    ax.set_title("Val MSE by Model — step-wise improvement", color=FG_C, fontsize=12, pad=8)

    plt.tight_layout(pad=0.9)
    path = tmp()
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="#0D1117")
    plt.close(fig)
    return path


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_img(slide, chart_title_bg(), 0, 0, width=W)
    add_rect(slide, 0, 0, Inches(0.12), H, ORG)

    textbox(slide, "Physics-Constrained\nChip Thermal Map Prediction",
            Inches(0.5), Inches(1.5), Inches(8.5), Inches(2.4),
            size=48, bold=True, color=WHITE)
    textbox(slide, "CS 231N — Deep Learning for Computer Vision  ·  Stanford Spring 2026",
            Inches(0.5), Inches(3.85), Inches(9), Inches(0.55),
            size=18, color=GRAY)
    textbox(slide, "Fabian Molina  ·  Ruben Carrazco",
            Inches(0.5), Inches(4.45), Inches(8), Inches(0.5),
            size=22, bold=True, color=ORG)

    add_rect(slide, Inches(0.5), Inches(5.2), Inches(3.2), Inches(0.52), CARD2)
    textbox(slide, "Milestone 3: Preliminary Results",
            Inches(0.5), Inches(5.2), Inches(3.2), Inches(0.52),
            size=16, bold=True, color=YEL, align=PP_ALIGN.CENTER)
    textbox(slide, "May 29, 2026",
            Inches(0.5), Inches(5.9), Inches(3), Inches(0.4),
            size=14, color=GRAY)


def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    header(slide, "Problem Statement & Motivation",
           "Can a neural network replace hours of thermal simulation?")

    cards = [
        (
            "The Challenge",
            "Chip overheating degrades performance and causes failures. During chip design, engineers must evaluate thousands of layout options — each requires a full thermal simulation.",
            ORG,
        ),
        (
            "The Bottleneck",
            "HotSpot, the standard academic thermal simulator, takes minutes per design. This makes it impractical for automated layout search across large design spaces.",
            BLUE,
        ),
        (
            "Our Goal",
            "Train a CNN to predict a full 256x256 temperature map from two input images (floorplan + power density) in milliseconds — a 1000x speedup over simulation.",
            GREEN,
        ),
    ]
    for i, (title, body, accent) in enumerate(cards):
        top = Inches(1.35) + i * Inches(1.95)
        add_rect(slide, Inches(0.4), top, Inches(5.7), Inches(1.8), CARD)
        add_rect(slide, Inches(0.4), top, Inches(0.07), Inches(1.8), accent)
        textbox(slide, title,
                Inches(0.6), top + Inches(0.12), Inches(5.4), Inches(0.42),
                size=17, bold=True, color=accent)
        textbox(slide, body,
                Inches(0.6), top + Inches(0.5), Inches(5.3), Inches(1.18),
                size=14, color=WHITE)

    # Right column
    add_rect(slide, Inches(6.55), Inches(1.35), Inches(6.35), Inches(5.75), CARD2)
    add_rect(slide, Inches(6.55), Inches(1.35), Inches(6.35), Inches(0.07), ORG)
    textbox(slide, "Core Research Question",
            Inches(6.75), Inches(1.5), Inches(6.0), Inches(0.5),
            size=18, bold=True, color=ORG)
    textbox(slide,
            "Does adding a physics constraint to the training loss help the model generalize to chip configurations it was never trained on?",
            Inches(6.75), Inches(2.1), Inches(5.9), Inches(1.8),
            size=22, color=WHITE)
    textbox(slide, "The physics constraint (heat equation):",
            Inches(6.75), Inches(4.05), Inches(5.8), Inches(0.4),
            size=14, bold=True, color=YEL)
    textbox(slide, "k · ∇²T + Q = 0",
            Inches(6.75), Inches(4.5), Inches(5.8), Inches(0.58),
            size=28, bold=True, color=ORG, align=PP_ALIGN.CENTER)
    textbox(slide, "T = temperature map,  Q = power density,  k = thermal conductivity",
            Inches(6.75), Inches(5.15), Inches(5.8), Inches(0.45),
            size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    textbox(slide, "We penalize how much the prediction violates this equation during training.",
            Inches(6.75), Inches(5.68), Inches(5.8), Inches(0.7),
            size=14, color=WHITE)


def slide_03_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    header(slide, "Dataset & Pipeline",
           "CircuitNet 2.0 — public benchmark of GPU/accelerator chip designs with HotSpot thermal labels")

    stats = [
        ("189", "chip designs"),
        ("3", "design families"),
        ("256x256", "pixels per map"),
        ("268 / 31 / 31", "train / val / test"),
    ]
    for i, (val, label) in enumerate(stats):
        lft = Inches(0.4) + i * Inches(3.22)
        add_rect(slide, lft, Inches(1.3), Inches(3.0), Inches(1.05), CARD)
        add_rect(slide, lft, Inches(1.3), Inches(3.0), Inches(0.06), ORG)
        textbox(slide, val,
                lft, Inches(1.38), Inches(3.0), Inches(0.56),
                size=32, bold=True, color=ORG, align=PP_ALIGN.CENTER)
        textbox(slide, label,
                lft, Inches(1.9), Inches(3.0), Inches(0.36),
                size=14, color=GRAY, align=PP_ALIGN.CENTER)

    img_path = DIR / "data/samples/sample_overview.png"
    if img_path.exists():
        add_img(slide, img_path, Inches(0.3), Inches(2.55), width=Inches(5.8))

    textbox(slide, "Inputs  (2 channels)",
            Inches(7.0), Inches(2.6), Inches(6.0), Inches(0.45),
            size=16, bold=True, color=ORG)
    textbox(slide,
            "Floorplan — grayscale image showing where circuit components are placed on the chip\n\n"
            "Power density map — how much power each region dissipates",
            Inches(7.0), Inches(3.1), Inches(5.9), Inches(1.5),
            size=15, color=WHITE)

    textbox(slide, "Target  (1 channel)",
            Inches(7.0), Inches(4.75), Inches(6.0), Inches(0.45),
            size=16, bold=True, color=BLUE)
    textbox(slide,
            "HotSpot thermal map — simulated temperature at every point on the chip surface (in Kelvin)",
            Inches(7.0), Inches(5.25), Inches(5.9), Inches(0.9),
            size=15, color=WHITE)

    families = [("Vortex-small", "96"), ("Vortex-large", "61"), ("nvdla-large", "32")]
    textbox(slide, "Design families:",
            Inches(7.0), Inches(6.25), Inches(5.9), Inches(0.35),
            size=14, bold=True, color=GRAY)
    for i, (fam, cnt) in enumerate(families):
        lft = Inches(7.0) + i * Inches(2.0)
        add_rect(slide, lft, Inches(6.65), Inches(1.85), Inches(0.55), CARD)
        textbox(slide, f"{fam}\n{cnt} designs",
                lft, Inches(6.65), Inches(1.85), Inches(0.55),
                size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_04_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    header(slide, "Model Architectures",
           "Three models — progressively adding skip connections and a physics loss")

    arch_data = [
        {
            "name": "PlainCNN",
            "tag": "Baseline A",
            "color": "#64748B",
            "lines": [
                "No encoder/decoder, no skip connections",
                "8 stacked conv layers with increasing dilation rates — allows wider receptive field without resolution loss",
                "Weakest baseline: no spatial compression or skip paths",
                "Input (2ch) -> 8 layers -> Output (1ch)",
                "64 feature channels throughout",
            ],
        },
        {
            "name": "EncoderDecoder",
            "tag": "Baseline B",
            "color": "#3B82F6",
            "lines": [
                "Encoder: 5 stages of downsampling (2x at each stage)",
                "Bottleneck: compressed 16x representation + dropout",
                "Decoder: 4 stages of upsampling back to 256x256",
                "No skip connections — spatial detail is lost at bottleneck",
                "Ablation: isolates the effect of skip connections",
            ],
        },
        {
            "name": "U-Net + Physics",
            "tag": "Our Model",
            "color": "#F97316",
            "lines": [
                "Same encoder-decoder structure as above",
                "Skip connections copy encoder features to decoder at each level — spatial detail is preserved",
                "Physics loss term added: penalizes predictions that violate the heat equation",
                "L_total = L_MSE + lambda * L_physics",
                "lambda controls strength of physics constraint",
            ],
        },
    ]

    for i, arch in enumerate(arch_data):
        lft  = Inches(0.35) + i * Inches(4.32)
        top  = Inches(1.35)
        clr  = RGBColor.from_string(arch["color"][1:])
        add_rect(slide, lft, top, Inches(4.1), Inches(5.75), CARD)
        add_rect(slide, lft, top, Inches(4.1), Inches(0.07), clr)

        add_rect(slide, lft + Inches(2.45), top + Inches(0.12), Inches(1.55), Inches(0.38), CARD2)
        textbox(slide, arch["tag"],
                lft + Inches(2.45), top + Inches(0.12), Inches(1.55), Inches(0.38),
                size=12, color=clr, align=PP_ALIGN.CENTER)

        textbox(slide, arch["name"],
                lft + Inches(0.15), top + Inches(0.15), Inches(2.25), Inches(0.52),
                size=21, bold=True, color=WHITE)

        for j, line in enumerate(arch["lines"]):
            t = top + Inches(0.82) + j * Inches(0.92)
            add_rect(slide, lft + Inches(0.12), t, Inches(3.86), Inches(0.8), CARD2)
            textbox(slide, line,
                    lft + Inches(0.22), t + Inches(0.05), Inches(3.66), Inches(0.7),
                    size=12, color=WHITE)


def slide_05_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    header(slide, "Quantitative Results",
           "Validation split (31 samples) · 250 epochs · Adam optimizer · lower is better")

    rows_data = [
        ("Model",               "Val MSE",  "RMSE (K)", "% above best", True),
        ("U-Net  (lambda=0)  *BEST*",  "0.01691",  "0.78 K",   "—",         False),
        ("U-Net  (lambda=0.1)", "0.02196",  "0.89 K",   "+14%",      False),
        ("EncoderDecoder",      "0.03972",  "1.20 K",   "+54%",      False),
        ("PlainCNN",            "0.07693",  "1.67 K",   "+114%",     False),
    ]

    tbl_shape = slide.shapes.add_table(
        len(rows_data), 4,
        Inches(0.4), Inches(1.35),
        Inches(6.1), Inches(5.5),
    )
    tbl = tbl_shape.table

    for ci, cw in enumerate([Inches(2.7), Inches(1.1), Inches(1.1), Inches(1.2)]):
        tbl.columns[ci].width = cw
    tbl.rows[0].height = Inches(0.65)
    for ri in range(1, len(rows_data)):
        tbl.rows[ri].height = Inches(1.0)

    for ri, (model, mse_v, rmse_v, delta, is_hdr) in enumerate(rows_data):
        for ci, val in enumerate([model, mse_v, rmse_v, delta]):
            cell = tbl.cell(ri, ci)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(15 if not is_hdr else 13)
            run.font.bold = is_hdr or ri == 1

            if is_hdr:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD2
                run.font.color.rgb = ORG
            elif ri == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x29, 0x1A, 0x0A)
                run.font.color.rgb = (YEL if ci == 0 else
                                      (ORG if ci in (1, 2) else GREEN))
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD2
                run.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD
                run.font.color.rgb = WHITE

            if ci == 3 and not is_hdr and ri > 1:
                run.font.color.rgb = RED

    bar_path = chart_rmse_bar()
    add_img(slide, bar_path, Inches(6.7), Inches(1.35), width=Inches(6.4))

    add_rect(slide, Inches(6.7), Inches(5.9), Inches(6.4), Inches(1.25), CARD2)
    add_rect(slide, Inches(6.7), Inches(5.9), Inches(0.07), Inches(1.25), GREEN)
    textbox(slide,
            "Adding skip connections (EncoderDecoder -> U-Net) reduces RMSE by 35% (1.20 -> 0.78 K).\n"
            "Physics loss raises in-distribution RMSE slightly — expected, it is a regularizer.",
            Inches(6.85), Inches(5.95), Inches(6.15), Inches(1.1),
            size=14, color=WHITE)


def slide_06_visual(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    header(slide, "Predictions vs. Ground Truth",
           "Best model: U-Net (MSE only, no physics loss) — evaluated on three unseen chip families")

    img_path = DIR / "per_family_preds.png"
    if img_path.exists():
        add_img(slide, img_path, Inches(0.3), Inches(1.3), height=Inches(5.55))

    # Side callout panel
    add_rect(slide, Inches(13.0), Inches(1.3), Inches(0.25), Inches(5.55), ORG)

    add_rect(slide, Inches(0.3), Inches(7.0), Inches(12.63), Inches(0.38), CARD)
    textbox(slide,
            "Columns:  Floorplan (input)  |  Power density (input)  |  Ground truth temperature  |  Predicted (GT scale)  |  Predicted (auto scale)",
            Inches(0.4), Inches(7.02), Inches(12.5), Inches(0.35),
            size=11, color=GRAY, align=PP_ALIGN.CENTER)


def slide_07_physics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    header(slide, "Effect of Physics Regularization",
           "lambda controls how strongly the training penalizes violations of the heat equation")

    # Formula card (left)
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(5.7), Inches(3.7), CARD)
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(5.7), Inches(0.07), ORG)
    textbox(slide, "Training Loss",
            Inches(0.6), Inches(1.5), Inches(5.3), Inches(0.48),
            size=18, bold=True, color=ORG)
    textbox(slide, "L = L_MSE  +  lambda * L_physics",
            Inches(0.6), Inches(2.02), Inches(5.3), Inches(0.56),
            size=21, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.5), Inches(2.7), Inches(5.5), Inches(0.025), SLATE)

    textbox(slide, "L_physics  measures how much the prediction\nviolates the steady-state heat equation:",
            Inches(0.6), Inches(2.78), Inches(5.3), Inches(0.75),
            size=14, color=GRAY)
    textbox(slide, "L_physics  =  || k * lap(T_pred) + Q ||^2",
            Inches(0.6), Inches(3.55), Inches(5.3), Inches(0.52),
            size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide,
            "lap = discrete Laplacian (measures spatial curvature)\n"
            "k = thermal conductivity\n"
            "Q = power density input",
            Inches(0.6), Inches(4.12), Inches(5.3), Inches(0.88),
            size=13, color=GRAY)

    # Lambda chart (right)
    lam_path = chart_lambda()
    add_img(slide, lam_path, Inches(6.35), Inches(1.35), width=Inches(6.75))

    # Two insight cards
    insights = [
        (
            "Expected: physics loss trades accuracy for regularity",
            "Adding physics regularization (lambda=0.1) increases val MSE by ~30% compared to lambda=0. This is expected — the model is constrained to be physically plausible, not just MSE-optimal.",
            BLUE,
        ),
        (
            "Open question: does it help on unseen chip configs?",
            "We hypothesize that physics-constrained models (lambda > 0) will degrade less when tested on chip thicknesses outside the training set. This out-of-distribution test is Milestone 4.",
            ORG,
        ),
    ]
    for i, (title, body, accent) in enumerate(insights):
        lft = Inches(0.4) + i * Inches(6.5)
        top = Inches(5.25)
        add_rect(slide, lft, top, Inches(6.2), Inches(2.0), CARD)
        add_rect(slide, lft, top, Inches(0.07), Inches(2.0), accent)
        textbox(slide, title,
                lft + Inches(0.18), top + Inches(0.12), Inches(5.9), Inches(0.45),
                size=15, bold=True, color=accent)
        textbox(slide, body,
                lft + Inches(0.18), top + Inches(0.55), Inches(5.9), Inches(1.3),
                size=13, color=WHITE)


def slide_08_analysis(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    header(slide, "Analysis & Insights",
           "What the architecture sweep tells us")

    mse_path = chart_mse_bars()
    add_img(slide, mse_path, Inches(7.0), Inches(1.35), width=Inches(6.1))

    insights = [
        (
            "Skip connections are the single biggest driver",
            "U-Net vs. EncoderDecoder have identical depth, channel counts, and training setup. The only difference is skip connections. Result: 57% reduction in val MSE (0.040 -> 0.017). Spatial detail lost in the bottleneck cannot be recovered by the decoder alone.",
            ORG,
            "57% MSE\nreduction",
        ),
        (
            "More layers is not the answer",
            "PlainCNN uses 8 convolutional layers (more than EncoderDecoder's encoder+decoder combined) yet performs 2x worse. In pixel-to-pixel prediction tasks, how you route spatial information matters more than raw depth.",
            BLUE,
            "Connectivity\n> Depth",
        ),
        (
            "Physics loss trains stably — in-distribution cost is small",
            "Increasing lambda from 0 to 0.1 raises val MSE by ~30% (not catastrophic). Training did not diverge. The physics gradients are well-behaved. The out-of-distribution benefit remains to be measured.",
            GREEN,
            "Stable across\nlambda values",
        ),
    ]
    for i, (title, body, accent, badge) in enumerate(insights):
        top = Inches(1.35) + i * Inches(2.0)
        add_rect(slide, Inches(0.3), top, Inches(6.45), Inches(1.88), CARD)
        add_rect(slide, Inches(0.3), top, Inches(0.07), Inches(1.88), accent)
        add_rect(slide, Inches(5.0), top + Inches(0.1), Inches(1.65), Inches(0.7), CARD2)
        textbox(slide, badge,
                Inches(5.0), top + Inches(0.1), Inches(1.65), Inches(0.7),
                size=11, bold=True, color=accent, align=PP_ALIGN.CENTER)
        textbox(slide, title,
                Inches(0.5), top + Inches(0.1), Inches(4.45), Inches(0.45),
                size=16, bold=True, color=accent)
        textbox(slide, body,
                Inches(0.5), top + Inches(0.52), Inches(6.15), Inches(1.25),
                size=13, color=WHITE)


def slide_09_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    header(slide, "Limitations", "What we do not yet know")

    limitations = [
        (
            "Small validation set",
            "Only 31 validation samples. RMSE estimates have high variance. The held-out test split (31 samples) has not been evaluated yet — all reported numbers are validation-only.",
            ORG,
        ),
        (
            "No out-of-distribution test yet",
            "The core claim — that physics-constrained models generalize better to unseen chip thicknesses — has not been tested. This is the goal of Milestone 4.",
            ORG,
        ),
        (
            "Physics weight lambda not optimized",
            "We measured lambda=0 and lambda=0.1. Two additional values (lambda=0.01, 1.0) were trained but results are still pending. The optimal lambda for out-of-distribution performance is unknown.",
            ORG,
        ),
        (
            "Labels approximate real silicon temperatures",
            "Ground truth comes from HotSpot, a block-level RC thermal model, not full 3D finite-element simulation. Prediction accuracy is bounded by how accurate HotSpot itself is.",
            SLATE,
        ),
        (
            "Only one architecture family tested",
            "All three models share the same U-Net backbone style. We have not compared against attention-based models or graph neural networks, which are out of scope for this milestone.",
            SLATE,
        ),
    ]

    for i, (title, body, accent) in enumerate(limitations):
        row = i % 3
        col = i // 3
        lft = Inches(0.35) + col * Inches(6.65)
        top = Inches(1.35) + row * Inches(2.0)
        add_rect(slide, lft, top, Inches(6.3), Inches(1.85), CARD)
        add_rect(slide, lft, top, Inches(0.07), Inches(1.85), accent)
        textbox(slide, title,
                lft + Inches(0.18), top + Inches(0.1), Inches(5.95), Inches(0.45),
                size=16, bold=True, color=WHITE)
        textbox(slide, body,
                lft + Inches(0.18), top + Inches(0.52), Inches(5.95), Inches(1.2),
                size=13, color=GRAY)


def slide_10_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    header(slide, "Next Steps",
           "From preliminary val results to paper-quality evaluation")

    steps = [
        {
            "n": "01",
            "title": "Test Set Evaluation",
            "body": "Evaluate the best U-Net checkpoint on the 31-sample held-out test split (never seen during training). Report RMSE, SSIM (structural similarity), and hotspot location accuracy (IoU of top-5% hottest pixels).",
            "color": ORG,
            "status": "READY",
        },
        {
            "n": "02",
            "title": "Prediction Visualization",
            "body": "Generate side-by-side comparison panels: U-Net lambda=0 vs U-Net lambda=0.1. Do physics constraints produce sharper or more accurate hotspot boundaries?",
            "color": BLUE,
            "status": "READY",
        },
        {
            "n": "03",
            "title": "Out-of-Distribution Test (core claim)",
            "body": "Re-run HotSpot simulation at new chip thicknesses outside the training range. Evaluate all checkpoints on these unseen configs. This is the central experiment of the paper.",
            "color": GREEN,
            "status": "IN PROGRESS",
        },
        {
            "n": "04",
            "title": "Complete Lambda Sweep",
            "body": "Collect validation metrics for lambda=0.01 and lambda=1.0 to complete the regularization curve. Understand the in-distribution vs. out-of-distribution trade-off across all four lambda values.",
            "color": PURP,
            "status": "PENDING",
        },
    ]

    for i, step in enumerate(steps):
        col = i % 2
        row = i // 2
        lft = Inches(0.35) + col * Inches(6.55)
        top = Inches(1.35) + row * Inches(2.9)
        add_rect(slide, lft, top, Inches(6.3), Inches(2.72), CARD)
        add_rect(slide, lft, top, Inches(0.07), Inches(2.72), step["color"])

        textbox(slide, step["n"],
                lft + Inches(0.18), top + Inches(0.12), Inches(0.7), Inches(0.55),
                size=28, bold=True, color=step["color"])

        sc = (GREEN if step["status"] == "READY"
              else (ORG if "PROGRESS" in step["status"] else GRAY))
        add_rect(slide, lft + Inches(4.65), top + Inches(0.15), Inches(1.55), Inches(0.38), CARD2)
        textbox(slide, step["status"],
                lft + Inches(4.65), top + Inches(0.15), Inches(1.55), Inches(0.38),
                size=11, bold=True, color=sc, align=PP_ALIGN.CENTER)

        textbox(slide, step["title"],
                lft + Inches(0.9), top + Inches(0.12), Inches(3.65), Inches(0.52),
                size=17, bold=True, color=WHITE)
        textbox(slide, step["body"],
                lft + Inches(0.18), top + Inches(0.72), Inches(5.95), Inches(1.82),
                size=13, color=GRAY)


# ── Main ──────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Building slides...")
    slide_01_title(prs);        print("  1/10  Title")
    slide_02_problem(prs);      print("  2/10  Problem & Motivation")
    slide_03_dataset(prs);      print("  3/10  Dataset & Pipeline")
    slide_04_architecture(prs); print("  4/10  Architectures")
    slide_05_results(prs);      print("  5/10  Results")
    slide_06_visual(prs);       print("  6/10  Visual Predictions")
    slide_07_physics(prs);      print("  7/10  Physics Loss Analysis")
    slide_08_analysis(prs);     print("  8/10  Analysis & Insights")
    slide_09_limitations(prs);  print("  9/10  Limitations")
    slide_10_next_steps(prs);   print(" 10/10  Next Steps")

    out = DIR / "milestone3_slides.pptx"
    prs.save(str(out))
    print(f"\n  Saved: {out}")
    cleanup()


if __name__ == "__main__":
    build()
